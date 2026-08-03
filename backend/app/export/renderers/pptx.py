import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.export.renderers.base import BaseRenderer
from app.export.renderers.common import iter_renderable_sections, build_section_blocks

logger = logging.getLogger("proposalcraft.export.pptx")

BLUE = RGBColor(26, 86, 219)
DARK = RGBColor(51, 51, 51)
GRAY = RGBColor(100, 100, 100)


class PPTXRenderer(BaseRenderer):
    extension = ".pptx"

    def render(self, proposal: dict, output_path: str) -> str:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        self._add_cover_slide(prs, proposal.get("metadata", {}))

        for key, title, data in iter_renderable_sections(proposal):
            self._add_section_slide(prs, title, build_section_blocks(key, data))

        path = Path(output_path).with_suffix(".pptx")
        prs.save(str(path))
        logger.info("PPTX exported to %s", path)
        return str(path)

    def _add_cover_slide(self, prs: Presentation, meta: dict):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 244, 255)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = meta.get("proposal_title", "Proposal")
        p.font.size = Pt(40)
        p.font.color.rgb = BLUE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        subtitle = meta.get("subtitle", "")
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.font.color.rgb = GRAY
            p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = ""
        p3.space_before = Pt(20)
        for label, key in [("Prepared for", "prepared_for"), ("Date", "date")]:
            val = meta.get(key)
            if val:
                p4 = tf.add_paragraph()
                p4.text = f"{label}: {val}"
                p4.font.size = Pt(14)
                p4.font.color.rgb = DARK
                p4.alignment = PP_ALIGN.CENTER

    def _add_section_slide(self, prs: Presentation, title: str, blocks: list[tuple]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = BLUE
        p.font.bold = True

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
        ctf = content_box.text_frame
        ctf.word_wrap = True

        first_para = True
        for block in blocks:
            kind = block[0]
            if kind == "h3":
                para = ctf.add_paragraph()
                para.text = block[1]
                para.font.size = Pt(20)
                para.font.color.rgb = BLUE
                para.font.bold = True
                para.space_before = Pt(8)
                para.space_after = Pt(2)
            elif kind == "p":
                para = ctf.add_paragraph()
                para.text = block[1]
                para.font.size = Pt(14)
                para.font.color.rgb = DARK
                para.space_after = Pt(4)
            elif kind == "bullets":
                for item in block[1]:
                    para = ctf.add_paragraph()
                    para.text = f"• {item}"
                    para.font.size = Pt(14)
                    para.font.color.rgb = DARK
                    para.space_after = Pt(2)
            elif kind == "table":
                headers, rows = block[1], block[2]
                para = ctf.add_paragraph()
                para.text = " | ".join(str(h) for h in headers)
                para.font.size = Pt(12)
                para.font.bold = True
                para.space_after = Pt(2)
                for row in rows:
                    para = ctf.add_paragraph()
                    para.text = " | ".join(str(v) for v in row)
                    para.font.size = Pt(12)
                    para.font.color.rgb = DARK
                    para.space_after = Pt(1)
            first_para = False

        if first_para:
            ctf.paragraphs[0].text = title
            ctf.paragraphs[0].font.size = Pt(24)
