"""Format parity: all four renderers must produce the same sections in the same
canonical order with the same module names — a single normalization step."""

import re

import pytest


def _fixture_proposal() -> dict:    return {
        "metadata": {
            "proposal_title": "Car Marketplace Proposal",
            "subtitle": "Web platform",
            "prepared_for": "Test Client",
            "prepared_by": "Test Agency",
            "date": "2026-08-03",
            "version": "1.0",
            "company_name": "Test Agency",
        },
        "about_company": {
            "who_we_are": "We build car marketplace platforms.",
            "experience": "Years of marketplace experience.",
            "why_choose_us": ["Fast delivery", "Senior team"],
        },
        "executive_summary": {
            "business_overview": "Client runs a peer-to-peer car marketplace.",
            "problem_statement": "Listing management is manual.",
            "proposed_solution": "An automated marketplace platform.",
            "expected_roi": "30% faster listings.",
            "business_value": "Scalable growth.",
            "key_benefits": ["Faster listings", "Better matching"],
        },
        "client_understanding": {
            "business_overview": "Car dealers and private sellers.",
            "current_challenges": ["Manual listing", "Slow search"],
            "business_goals": ["Increase listings", "Improve search"],
        },
        "requirement_analysis": {
            "functional_requirements": ["Car listing management", "Search and filter"],
            "non_functional_requirements": ["Fast search", "High availability"],
            "assumptions": ["Sellers provide photos"],
        },
        "proposed_solution": {
            "overview": "End-to-end marketplace platform.",
            "architecture": "Modular architecture.",
            "workflow": "List to sell in minutes.",
            "security": "Role-based access.",
            "scalability": "Horizontal scaling.",
        },
        "module_breakdown": {
            "standard_modules": [
                {"module_name": "Car Listing Management", "description": "Manage listings"},
                {"module_name": "Search and Filter", "description": "Find cars"},
            ],
            "custom_modules": [
                {"module_name": "Test Drive Booking", "description": "Book test drives"},
            ],
        },
        "technology_stack": {
            "frontend": ["React 19", "TypeScript"],
            "backend": ["FastAPI"],
            "database": ["MongoDB"],
            "rationale": "React powers the listing UI while FastAPI serves the API.",
        },
        "methodology": {
            "approach": "Agile-Scrum",
            "phases": ["Discovery & Architecture (3 weeks) — Stakeholder interviews"],
            "ceremonies": ["Daily standups"],
        },
        "timeline": {
            "phases": [
                {"name": "Discovery", "duration_weeks": 3, "activities": ["Interviews", "Analysis"]},
            ],
            "milestones": ["M1: Design freeze"],
        },
        "deliverables": ["Car Listing Management"],
        "pricing": {
            "one_time_cost": 15000,
            "monthly_cost": 1500,
            "annual_cost": 33000,
            "five_year_tco": 105000,
            "total_effort_hours": 240,
            "support_hours_included": 30,
            "effort_breakdown": {
                "Car Listing Management": {"hours": 80, "hourly_rate": 72, "cost": 5760},
                "Search and Filter": {"hours": 80, "hourly_rate": 72, "cost": 5760},
            },
            "payment_options": [
                {"type": "Milestone Based", "description": "30/40/30", "amount": 15000},
            ],
        },
        "custom_development_charges": [
            {"name": "Extra custom work", "cost": 2000},
        ],
        "sla": {
            "uptime_guarantee": "99.9%",
            "sla_tiers": [{"priority": "P1", "description": "Critical", "response_time": "1h", "resolution_time": "4h"}],
        },
        "support": {
            "recommended_plan": "Standard",
            "hours": "30 hours/month",
            "response_time": "4 hours",
            "channels": "Email, Chat",
            "included": ["Bug fixes"],
        },
        "security": {
            "authentication": "OAuth 2.0 with MFA.",
            "authorization": "Role-based access control.",
            "encryption": "AES-256 at rest.",
            "audit_logs": "Full audit trail.",
            "backup": "Daily backups.",
        },
        "terms": {
            "assumptions": ["Client provides APIs."],
            "exclusions": ["Third-party licenses."],
            "confidentiality": "NDA applies.",
            "warranty": "90-day warranty.",
        },
        "team": [
            {"role": "Project Manager", "seniority": "Senior", "allocation": "full_time"},
            {"role": "Backend Developer", "seniority": "Senior", "allocation": "full_time"},
        ],
        "conclusion": {
            "summary": "Ready to build.",
            "next_steps": ["Sign agreement", "Kickoff"],
        },
    }


def _section_titles(content: str) -> list[str]:
    if content.lstrip().startswith("<?xml") or content.startswith("PK"):
        return []
    titles = re.findall(r"(?:<h1>|<h2>|Heading level 1|(?m)^[A-Z][A-Za-z &]+$)", content)
    return titles


# Canonical order from SECTION_RULES (section_rules.py) — single source of truth.
_EXPECTED_TITLES = [
    "About Us", "Executive Summary", "Client Understanding", "Proposed Solution",
    "Requirement Analysis", "Module Breakdown", "Technology Stack", "Methodology",
    "Timeline", "Deliverables", "Investment", "Custom Development Charges",
    "Service Level Agreement", "Support Plan", "Security", "Terms & Conditions",
    "Team", "Conclusion",
]


class TestRendererParity:
    @pytest.fixture(scope="class")
    def outputs(self):
        from app.export.service import export_service

        proposal = _fixture_proposal()
        results = export_service.export_all(proposal)
        assert all(isinstance(v, str) and v.endswith((".pdf", ".docx", ".pptx", ".html")) for v in results.values()), results
        return results

    def test_all_formats_succeed(self, outputs):
        assert set(outputs.keys()) == {"html", "pdf", "docx", "pptx"}

    def test_docx_contains_same_section_titles(self, outputs):
        from docx import Document

        doc = Document(outputs["docx"])
        titles = [p.text for p in doc.paragraphs if p.style.name.startswith("Heading 1")]
        assert titles == _EXPECTED_TITLES

    def test_docx_contains_same_module_names(self, outputs):
        from docx import Document

        doc = Document(outputs["docx"])
        text = "\n".join(p.text for p in doc.paragraphs)
        table_text = "\n".join(
            cell.text for t in doc.tables for row in t.rows for cell in row.cells
        )
        combined = text + "\n" + table_text
        for module in ["Car Listing Management", "Search and Filter", "Test Drive Booking"]:
            assert module in combined

    def test_pptx_contains_same_section_titles(self, outputs):
        from pptx import Presentation

        prs = Presentation(outputs["pptx"])
        titles = []
        for slide_idx, slide in enumerate(prs.slides):
            if slide_idx == 0:
                continue  # cover slide
            top_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
            if not top_shapes:
                continue
            title_shape = min(top_shapes, key=lambda s: s.top)
            titles.append(title_shape.text_frame.text.strip())
        assert titles == _EXPECTED_TITLES

    def test_pptx_contains_same_module_names(self, outputs):
        from pptx import Presentation

        prs = Presentation(outputs["pptx"])
        text = "\n".join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        for module in ["Car Listing Management", "Search and Filter", "Test Drive Booking"]:
            assert module in text

    def test_html_contains_same_section_titles(self, outputs):
        content = open(outputs["html"], encoding="utf-8").read()
        titles = re.findall(r"<h2>([^<]+)</h2>", content)
        titles = [t.replace("&amp;", "&") for t in titles]
        assert titles[1:] == _EXPECTED_TITLES  # [0] is the cover-page subtitle

    def test_html_contains_same_module_names(self, outputs):
        content = open(outputs["html"], encoding="utf-8").read()
        for module in ["Car Listing Management", "Search and Filter", "Test Drive Booking"]:
            assert module in content

    def test_pdf_text_contains_same_section_titles(self, outputs):
        try:
            from pdfminer.high_level import extract_text
        except OSError:
            pytest.skip("pdfminer not available")
        text = extract_text(outputs["pdf"])
        for title in ["About Us", "Executive Summary", "Module Breakdown",
                      "Technology Stack", "Investment", "Security", "Team", "Conclusion"]:
            assert title in text

    def test_pdf_text_contains_same_module_names(self, outputs):
        try:
            from pdfminer.high_level import extract_text
        except OSError:
            pytest.skip("pdfminer not available")
        text = extract_text(outputs["pdf"])
        for module in ["Car Listing Management", "Search and Filter", "Test Drive Booking"]:
            assert module in text
